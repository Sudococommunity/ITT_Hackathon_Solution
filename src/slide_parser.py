"""
Slide parser for extracting structured project data from PPTX files.
Groups slides by project and extracts metadata.
"""
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches


def extract_slide_text(slide) -> str:
    """Extract all text from a single slide, preserving structure."""
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(shape.text.strip())
    return "\n".join(texts)


def extract_table_data(slide) -> list[dict]:
    """Extract data from tables in a slide."""
    tables = []
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            rows = []
            for row in table.rows:
                row_data = [cell.text.strip() for cell in row.cells]
                rows.append(row_data)
            tables.append(rows)
    return tables


def parse_project_metadata(text: str) -> dict:
    """Parse structured metadata from slide text."""
    metadata = {}

    # Extract project ID (e.g., P001, P002)
    pid_match = re.search(r'\b(P\d{3})\b', text)
    if pid_match:
        metadata['project_id'] = pid_match.group(1)

    # Extract domain
    domain_match = re.search(r'Domain\s*\n?\s*(.+?)(?:\n|Team Size|Duration|Type|$)', text, re.IGNORECASE)
    if domain_match:
        metadata['domain'] = domain_match.group(1).strip()

    # Extract team size
    team_match = re.search(r'Team Size\s*\n?\s*(\d+)\s*(?:people)?', text, re.IGNORECASE)
    if team_match:
        metadata['team_size'] = int(team_match.group(1))

    # Extract duration
    duration_match = re.search(r'Duration\s*\n?\s*(\d+)\s*(months?|weeks?)', text, re.IGNORECASE)
    if duration_match:
        metadata['duration_months'] = int(duration_match.group(1))
        metadata['duration_unit'] = duration_match.group(2).lower()

    # Extract type/category
    type_match = re.search(r'Type\s*\n?\s*(.+?)(?:\n|$)', text, re.IGNORECASE)
    if type_match:
        metadata['project_type'] = type_match.group(1).strip()

    # Extract technology stack
    tech_match = re.search(r'TECHNOLOGY STACK\s*\n?\s*(.+?)(?:\n(?:EXPECTED|KEY)|$)', text, re.DOTALL | re.IGNORECASE)
    if tech_match:
        raw_tech = tech_match.group(1).strip()
        techs = [t.strip() for t in re.split(r'\s*[·,]\s*|\n', raw_tech) if t.strip()]
        metadata['technologies'] = techs

    # Extract project title
    lines = text.split('\n')
    for line in lines:
        line = line.strip()
        # Skip the project ID line, look for the title
        if line and not re.match(r'^P\d{3}$', line) and not line.startswith('PROBLEM') and not line.startswith('SOLUTION') and not line.startswith('PROJECT') and not line.startswith('TECHNOLOGY') and not line.startswith('EXPECTED'):
            pid = metadata.get('project_id', '')
            # Title is usually the second meaningful line or after project ID
            if pid and line.startswith(pid):
                # Lines like "P001  ›  Customer Churn..."
                title_match = re.search(r'P\d{3}\s*›?\s*(.+?)(?:\s*—|$)', line)
                if title_match:
                    metadata['title'] = title_match.group(1).strip()
                    break
            elif len(line) > 5 and len(line) < 200 and not any(kw in line.upper() for kw in ['PROBLEM', 'SOLUTION', 'PROJECT DETAILS', 'TECHNOLOGY', 'EXPECTED', 'FULL TECHNOLOGY', 'KEY METRICS']):
                if 'title' not in metadata:
                    metadata['title'] = line
                    break

    return metadata


def parse_pptx(pptx_path: str) -> list[dict]:
    """
    Parse PPTX and return list of project records.
    Each project groups its slides (typically 2: summary + deep dive).
    """
    prs = Presentation(pptx_path)
    slides_data = []

    slide_num = 0
    for slide in prs.slides:
        slide_num += 1
        text = extract_slide_text(slide)
        tables = extract_table_data(slide)
        slides_data.append({
            'slide_number': slide_num,
            'text': text,
            'tables': tables,
        })

    # Group slides by project
    # Pattern: odd slides are summaries (P001), even slides are deep dives (P001 › ...)
    # First slide is the title/cover slide
    projects = []
    i = 0

    # Skip cover slide
    if slides_data and 'Project Repository' in slides_data[0]['text']:
        i = 1

    while i < len(slides_data):
        slide = slides_data[i]
        text = slide['text']

        # Check if this is a project summary slide (has PXXX pattern)
        pid_match = re.search(r'\b(P\d{3})\b', text)
        if pid_match:
            project_id = pid_match.group(1)
            metadata = parse_project_metadata(text)

            # Collect all slides for this project
            project_slides = [slide]
            slide_numbers = [slide['slide_number']]

            # Check next slide(s) for deep dive of same project
            j = i + 1
            while j < len(slides_data):
                next_text = slides_data[j]['text']
                if project_id in next_text and ('Deep Dive' in next_text or 'Technical' in next_text or 'FULL TECHNOLOGY' in next_text):
                    project_slides.append(slides_data[j])
                    slide_numbers.append(slides_data[j]['slide_number'])
                    j += 1
                else:
                    break

            # Merge all text from project slides
            full_text = "\n\n".join([s['text'] for s in project_slides])

            # Re-parse metadata from full text for completeness
            full_metadata = parse_project_metadata(full_text)
            metadata.update({k: v for k, v in full_metadata.items() if k not in metadata or not metadata[k]})

            projects.append({
                'project_id': metadata.get('project_id', project_id),
                'title': metadata.get('title', 'Unknown'),
                'domain': metadata.get('domain', 'Unknown'),
                'team_size': metadata.get('team_size', 0),
                'duration_months': metadata.get('duration_months', 0),
                'project_type': metadata.get('project_type', 'Unknown'),
                'technologies': metadata.get('technologies', []),
                'slide_numbers': slide_numbers,
                'full_text': full_text,
                'summary_text': slide['text'],
            })

            i = j
        else:
            i += 1

    return projects


if __name__ == "__main__":
    import sys
    sys.stdout.reconfigure(encoding='utf-8')
    projects = parse_pptx("D:/hackathone2/Dataset_project_repository.pptx")
    print(f"Parsed {len(projects)} projects")
    for p in projects[:3]:
        print(f"\n{p['project_id']}: {p['title']}")
        print(f"  Domain: {p['domain']}, Team: {p['team_size']}, Duration: {p['duration_months']}m")
        print(f"  Tech: {', '.join(p['technologies'][:5])}")
        print(f"  Slides: {p['slide_numbers']}")
